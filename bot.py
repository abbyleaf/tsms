# pip install python-telegram-bot
import telegram.bot
from telegram import InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (Updater, CommandHandler, MessageHandler,
                          Filters, ConversationHandler, CallbackQueryHandler)
import logging
from functools import wraps
import os
import re
from cred import bottoken, port, rtlurl, rtlheaders
from templates import login, welcome, examples, defaultbook
import json
import requests
from io import BytesIO

# pip install Unidecode
from unidecode import unidecode

# pip install fuzzywuzzy[speedup]
from fuzzywuzzy import fuzz

# pip install python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

# pip install requests
from requests import get
# pip install pyopenssl
ip = get('https://api.ipify.org').text
try:
    certfile = open("cert.pem")
    keyfile = open("private.key")
    certfile.close()
    keyfile.close()
except IOError:
    from OpenSSL import crypto
    key = crypto.PKey()
    key.generate_key(crypto.TYPE_RSA, 2048)
    cert = crypto.X509()
    cert.get_subject().CN = ip
    cert.set_serial_number(1000)
    cert.gmtime_adj_notBefore(0)
    cert.gmtime_adj_notAfter(10*365*24*60*60)
    cert.set_issuer(cert.get_subject())
    cert.set_pubkey(key)
    cert.sign(key, 'sha256')
    with open("cert.pem", "wt") as certfile:
        certfile.write(crypto.dump_certificate(
            crypto.FILETYPE_PEM, cert).decode('ascii'))
    with open("private.key", "wt") as keyfile:
        keyfile.write(crypto.dump_privatekey(
            crypto.FILETYPE_PEM, key).decode('ascii'))

logging.basicConfig(filename='debug.log', filemode='a+', format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
                    level=logging.DEBUG)
logger = logging.getLogger(__name__)


def verifieduser(func):
    @wraps(func)
    def wrapped(update, context, *args, **kwargs):
        user_id = str(update.effective_user.id)
        if user_id in users:
            return func(update, context, *args, **kwargs)
        else:
            update.message.reply_text("Press /start to continue")
    return wrapped


def loader():
    global users
    try:
        with open('users.json') as userfile:
            users = json.load(userfile)
    except:
        with open('users.json', 'w+') as userfile:
            users = {}

    global songs
    songs = {}
    global titles
    titles = {}
    number = re.compile('^\d')
    for filename in (x.name for x in os.scandir('./books') if x.is_file()):
        filepath = './books/' + filename
        with open(filepath, 'r', encoding='UTF8') as book:
            filename = filename.split('.')
            bookname = filename[0]
            print("Loading", bookname)
            song_number = None
            song_lyrics = ''
            for line in book:
                if number.search(line):
                    if song_number != None:
                        song_lyrics = song_lyrics.strip()
                        songs[song_number] = song_lyrics
                        song_lyrics = ''
                    song_lyrics += bookname + ' ' + line
                    line = line.split()
                    song_number = bookname + ' ' + str(line.pop(0))
                    song_title = ' '.join(line)
                    titles[song_number] = song_title
                else:
                    song_lyrics += line
            song_lyrics = song_lyrics.strip()
            songs[song_number] = song_lyrics

    global chords
    chords = {}
    number = re.compile('^\d')
    with open('./media/tsms_chords.txt', 'r', encoding='UTF8') as tsms_chords_file:
        print("Loading Chords")
        song_number = None
        song_lyrics = ''
        for line in tsms_chords_file:
            if number.search(line):
                if song_number != None:
                    song_lyrics = song_lyrics.strip()
                    chords[song_number] = song_lyrics
                    song_lyrics = ''
                song_lyrics += 'TSMS ' + line
                line = line.split()
                song_number = 'TSMS ' + str(line.pop(0))
            else:
                song_lyrics += line
        song_lyrics = song_lyrics.strip()
        chords[song_number] = song_lyrics

    global scores
    scores = {}
    with open('./media/scores.txt', 'r', encoding='UTF8') as scores_file:
        print("Loading Scores")
        for line in scores_file:
            line = line.strip()
            line = line.split('@')
            reference = line[1]
            line = line[0]
            line = line.split('_')
            number = line[0]
            scores.setdefault(number, []).append(reference)

    global mp3
    mp3 = {}
    with open('./media/mp3.txt', 'r', encoding='UTF8') as mp3_file:
        print("Loading MP3")
        for line in mp3_file:
            line = line.strip()
            line = line.split('@')
            reference = line[1]
            line = line[0]
            line = line.split('_')
            number = line[0]
            mp3.setdefault(number, []).append(reference)

    global videos
    videos = {}
    with open('./media/videos.txt', 'r', encoding='UTF8') as videos_file:
        print("Loading Videos")
        for line in videos_file:
            line = line.strip()
            line = line.split('@')
            reference = line[1]
            number = line[0]
            videos.setdefault(number, []).append(reference)

    global piano
    piano = {}
    with open('./media/wilds_piano.txt', 'r', encoding='UTF8') as piano_file:
        print("Loading Piano")
        for line in piano_file:
            line = line.strip()
            line = line.split('@')
            reference = line[1]
            number = line[0]
            videos.setdefault(number, []).append(reference)


def start(update, context):
    contact_keyboard = telegram.KeyboardButton(
        text="REGISTER", request_contact=True)
    reply_markup = telegram.ReplyKeyboardMarkup([[contact_keyboard]])
    update.message.reply_text(
        login, reply_markup=reply_markup, parse_mode=telegram.ParseMode.MARKDOWN)


def contact(update, context):
    user_id = str(update.effective_user.id)
    contact = update.message.contact
    phone = contact.phone_number
    if user_id != str(contact.user_id):
        update.message.reply_text(
            "Verification failed. Your ID does not match.", reply_markup=telegram.ReplyKeyboardRemove())
    elif phone.startswith('+65') or phone.startswith('65') or phone.startswith('+61') or phone.startswith('61'):
        first_name = contact.first_name
        last_name = contact.last_name
        full_name = (str(first_name or '') + ' ' +
                     str(last_name or '')).strip()
        users[user_id] = {'phone': phone, 'name': full_name, 'rtl': False}
        with open('users.json', 'w') as userfile:
            json.dump(users, userfile)
        update.message.reply_text(welcome, parse_mode=telegram.ParseMode.MARKDOWN,
                                  reply_markup=telegram.ReplyKeyboardRemove())
    else:
        update.message.reply_text(
            "Sorry, you are not permitted to enter at this time.", reply_markup=telegram.ReplyKeyboardRemove())


@verifieduser
def helptext(update, context):
    update.message.reply_text(examples, parse_mode=telegram.ParseMode.MARKDOWN)


@verifieduser
def go(update, context):
    message = update.message.text
    if re.search("[\u4e00-\u9FFF]", message):
        update.message.reply_text(
            "_Error: Chinese characters not supported. Search using hanyu pinyin instead._", parse_mode=telegram.ParseMode.MARKDOWN)
        return
    message = message.strip().upper()
    if message.isnumeric():
        message = int(message)
        message = defaultbook + ' ' + str(message)
    if message in songs:
        reply_header = message
        reply = songs[reply_header]
    else:
        context.bot.send_chat_action(
            chat_id=update.message.chat_id, action=telegram.ChatAction.TYPING)
        alpha = re.compile('[^a-zA-Z ]')
        message = unidecode(message)
        message = alpha.sub('', message)
        results = []
        titleresults = set()
        titlefound = None
        for number, title in titles.items():
            title = unidecode(title)
            title = alpha.sub('', title)
            title = title.upper()
            if message == title:
                titlefound = number
                reply_header = number
                reply = songs[number]
                if number.startswith(defaultbook):
                    break
            similarity = fuzz.partial_ratio(title, message)
            if similarity > 75:
                results.append((number, -similarity))
                titleresults.add(number)
        if not titlefound:
            for number, lyrics in songs.items():
                lyrics = unidecode(lyrics)
                lyrics = alpha.sub('', lyrics)
                lyrics = lyrics.upper()
                similarity = fuzz.partial_ratio(lyrics, message)
                if similarity > 75 and number not in titleresults:
                    results.append((number, -similarity))
            results = sorted(results, key=lambda x: (x[1], x[0]))
            if len(results) == 0:
                update.message.reply_text(
                    "_No matches found_", parse_mode=telegram.ParseMode.MARKDOWN)
                return
            search_return = "<u>SONG SEARCH</u> (in order of relevance)\n\n"
            countresult = 0
            for result in results:
                countresult += 1
                if countresult > 50:
                    break
                number = result[0]
                search_return += '<b>{}</b> {}\n'.format(
                    number, titles[number])
            if countresult < 50:
                search_return += '\n<i>{} results</i>'.format(
                    str(len(results)))
            else:
                search_return += '\n<i>{} results (showing top 50)</i>'.format(
                    str(len(results)))
            update.message.reply_text(
                search_return, parse_mode=telegram.ParseMode.HTML)
            return
    keyboard = []
    if reply_header in chords:
        data = 'CHORDS {}'.format(reply_header)
        keyboard.append([InlineKeyboardButton(
            "Guitar Chords", callback_data=data)])
    if reply_header in scores:
        data = 'SCORE {}'.format(reply_header)
        keyboard.append([InlineKeyboardButton(
            "Piano Score", callback_data=data)])
    if reply_header in mp3:
        data = 'MP3 {}'.format(reply_header)
        keyboard.append([InlineKeyboardButton(
            "MIDI Soundtrack", callback_data=data)])
    if titles[reply_header] in piano:
        data = 'PIANO {}'.format(titles[reply_header])
        keyboard.append([InlineKeyboardButton(
            "Piano Recording (Wilds)", callback_data=data)])
    if reply_header in videos:
        data = 'VIDEO {}'.format(reply_header)
        keyboard.append([InlineKeyboardButton(
            "Choir Recording (Lyric Video)", callback_data=data)])
    if '\n\n' in reply:
        data = 'PPT {}'.format(reply_header)
        keyboard.append([InlineKeyboardButton(
            "Powerpoint (Beta)", callback_data=data)])
    reply_markup = InlineKeyboardMarkup(keyboard)
    update.message.reply_text(
        reply, reply_markup=reply_markup, parse_mode=telegram.ParseMode.MARKDOWN, disable_web_page_preview=True)


def callbackquery(update, context):
    query = update.callback_query
    data = query.data
    errors = False
    if data.startswith('CHORDS '):
        data = data.replace('CHORDS ', '')
        context.bot.send_message(chat_id=query.message.chat_id,
                                 text=chords[data], parse_mode=telegram.ParseMode.MARKDOWN, disable_web_page_preview=True)
    elif data.startswith('SCORE '):
        data = data.replace('SCORE ', '')
        counter = 1
        for i in range(len(scores[data])):
            reference = scores[data][i]
            title = data + ' ' + titles[data]
            if len(scores[data]) > 1:
                title += ' ' + str(counter)
            try:
                context.bot.send_photo(
                    chat_id=query.message.chat_id, photo=reference, caption=title)
            except:
                errors = True
            counter += 1
    elif data.startswith('MP3'):
        data = data.replace('MP3 ', '')
        counter = 1
        for i in range(len(mp3[data])):
            reference = mp3[data][i]
            title = title = data + ' ' + titles[data]
            if len(mp3[data]) > 1:
                title += ' ' + str(counter)
            try:
                context.bot.send_audio(
                    chat_id=query.message.chat_id, audio=reference, caption=title)
            except:
                errors = True
            counter += 1
    elif data.startswith('PIANO'):
        data = data.replace('PIANO ', '')
        reference = piano[data]
        try:
            context.bot.send_audio(
                chat_id=query.message.chat_id, audio=reference, caption=data)
        except:
            errors = True
    elif data.startswith('VIDEO'):
        data = data.replace('VIDEO ', '')
        context.bot.send_message(chat_id=query.message.chat_id,
                                 text='_Connecting to Life RTL..._', parse_mode=telegram.ParseMode.MARKDOWN)
        if users[str(query.message.chat_id)]['rtl'] == True:
            try:
                count = len(videos[data])
                if count > 1:
                    context.bot.send_message(
                        chat_id=query.message.chat_id, text='_{} versions found. Sending..._'.format(str(count)), parse_mode=telegram.ParseMode.MARKDOWN)
                for video in videos[data]:
                    url = rtlurl + video
                    videofile = requests.get(
                        url, headers=rtlheaders, timeout=10).content
                    context.bot.send_video(
                        chat_id=query.message.chat_id, video=BytesIO(videofile))
            except:
                errors = True
        else:
            context.bot.send_message(
                chat_id=query.message.chat_id, text='_You are not authorised to use this feature._', parse_mode=telegram.ParseMode.MARKDOWN)
    elif data.startswith('PPT'):
        data = data.replace('PPT ', '')
        slides = makeppt(data)
        try:
            context.bot.send_document(
                chat_id=query.message.chat_id, document=slides)
        except:
            errors = True
    if errors:
        context.bot.answer_callback_query(
            query.id, text='An Error Occurred: one or more files were unable to send. Please notify an admin.', show_alert=True)
    else:
        context.bot.answer_callback_query(query.id)


def makeppt(data):
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    title = titles[data]
    text = songs[data]
    text = text.split('\n\n')
    text.pop(0)
    text = list(filter(None, text))

    originallen = len(text)
    chorus = None
    for i in range(originallen):
        stanza = text[i]
        if stanza.startswith("Chorus:") or stanza.startswith("Refrain:"):
            chorus = i
            break
    if chorus:
        i = chorus+2
        while True:
            text.insert(i, stanza)
            i += 2
            if i > len(text):
                break

    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    txBox = slide.shapes.add_textbox(0, 0, Inches(16), Inches(9))
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.add_paragraph()
    p.text = title + '\n(' + data + ')'
    p.font.size = Pt(60)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    for i in range(len(text)):
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)

        txBox = slide.shapes.add_textbox(Inches(15), 0, Inches(1), Inches(1))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = '{}/{}'.format(i+1, len(text))
        p.font.size = Pt(32)

        txBox = slide.shapes.add_textbox(0, 0, Inches(16), Inches(9))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        p = tf.add_paragraph()
        p.text = text[i]
        p.font.size = Pt(48)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

    pptxfile = BytesIO()
    pptxfile.name = data + '.pptx'
    prs.save(pptxfile)
    pptxfile.seek(0)
    return pptxfile


def main():
    updater = Updater(token=bottoken, workers=32, use_context=True)
    dp = updater.dispatcher

    dp.add_handler(CommandHandler("start", start))
    dp.add_handler(CommandHandler("help", helptext))
    dp.add_handler(MessageHandler(Filters.contact, contact))
    dp.add_handler(MessageHandler(Filters.text, go))
    dp.add_handler(CallbackQueryHandler(callbackquery))

    loader()

    # updater.start_polling()
    updater.start_webhook(listen='0.0.0.0',
                          port=port,
                          url_path=bottoken,
                          key='private.key',
                          cert='cert.pem',
                          webhook_url='https://{}:{}/{}'.format(ip, port, bottoken))

    print("Bot is running. Press Ctrl+C to stop.")
    updater.idle()
    print("Bot stopped.")


if __name__ == '__main__':
    main()
